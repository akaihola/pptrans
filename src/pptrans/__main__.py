"""Translate PowerPoint presentations."""

from __future__ import annotations

import shutil
from pathlib import Path
from typing import TYPE_CHECKING, Any

import click
import llm
from pptx import Presentation

from pptrans.cache import (
    commit_pending_cache_updates,
    generate_page_hash,
    load_cache,
    prepare_slide_for_translation,
    update_data_from_llm_response,
)
from pptrans.page_range import parse_page_range

if TYPE_CHECKING:
    from pptx.slide import Slide as PptxSlide  # pragma: no cover


EOL_MARKER = "<"


def _emit_save_message(mode: str, output_path: Path) -> None:
    """Emit the final save message to the console."""
    click.echo(f"Presentation saved in '{mode}' mode to: {output_path}")


def reverse_individual_words(text_string_with_eol: str) -> str:
    """Reverse words in a space-separated string, preserving an EOL_MARKER."""
    text_to_reverse = text_string_with_eol
    had_eol = False
    if text_string_with_eol.endswith(EOL_MARKER):
        text_to_reverse = text_string_with_eol[: -len(EOL_MARKER)]
        had_eol = True

    words = text_to_reverse.split(" ")
    reversed_words = [word[::-1] for word in words]
    result = " ".join(reversed_words)

    if had_eol:
        return result + EOL_MARKER
    return result


# Helper function for slide selection
def _handle_slide_selection(
    pages_option: str | None, num_original_slides: int
) -> set[int]:
    """Parse page range string and return a set of 0-indexed page numbers."""
    if pages_option:
        # parse_page_range is imported from pptrans.page_range
        return parse_page_range(pages_option, num_original_slides)
    return set(range(num_original_slides))


# Helper function to extract text runs from a slide
def _extract_run_info_from_slide(slide_obj: PptxSlide) -> list[dict[str, Any]]:
    """Extract all text run objects and their original text from a slide."""
    run_info_list: list[dict[str, Any]] = []
    for shape in slide_obj.shapes:
        if shape.has_text_frame:
            run_info_list.extend(
                {"original_text": run.text, "run_object": run}
                for paragraph in shape.text_frame.paragraphs
                for run in paragraph.runs
                if run.text
            )
        if shape.has_table:
            run_info_list.extend(
                {"original_text": run.text, "run_object": run}
                for row in shape.table.rows
                for cell in row.cells
                for paragraph in cell.text_frame.paragraphs
                for run in paragraph.runs
                if run.text
            )
    return run_info_list


# Helper function to build LLM prompt
def _build_llm_prompt_and_data(texts_for_llm: list[dict[str, Any]]) -> tuple[str, str]:
    """Construct the LLM prompt and the formatted data string."""
    formatted_text_for_llm = "\n".join(
        [f"{item['id']}:{item['text_to_send']}" for item in texts_for_llm]
    )
    # EOL_MARKER is a global constant in this module
    prompt_text = (
        "You are an expert Finnish to English translator. "
        "Translate the following text segments accurately from Finnish to "
        "English. "
        "Each segment is prefixed with a unique ID (e.g., pg1_txt0, pg1_txt1). "
        "IMPORTANT: A sequence of text items (e.g., pg1_txt0, pg1_txt1, "
        "pg1_txt2) may represent a single continuous sentence that has been "
        "split due to formatting. Interpret and translate such sequences as a "
        "coherent whole sentence to maintain context and flow. "
        f"The text for each ID might end with an EOL marker: '{EOL_MARKER}'. "
        "Your response MUST consist ONLY of the translated segments, each "
        "prefixed with its original ID, "
        "and each on a new line. Maintain the exact ID and format. "
        "PRESERVE ALL LEADING AND TRAILING WHITESPACE from the original "
        "segment in your translation. "
        f"If an EOL marker '{EOL_MARKER}' was present at the end of the input segment, "
        f"IT MUST be present at the end of your translated segment, including "
        "any whitespace before it.\n"
        "For example, if you receive:\n"
        "pg1_txt0: Tämä on pitkä \n"
        "pg1_txt1:lause, joka on \n"
        f"pg1_txt2:jaettu.{EOL_MARKER}\n"
        f"pg1_txt3:    Toinen lause.   {EOL_MARKER}\n"
        "pg2_txt0: Yksittäinen.\n"
        "You MUST return:\n"
        "pg1_txt0: This is a long \n"
        "pg1_txt1:sentence that has been \n"
        f"pg1_txt2:split.{EOL_MARKER}\n"
        f"pg1_txt3:    Another sentence.   {EOL_MARKER}\n"
        "pg2_txt0: Standalone.\n\n"
        "Do not add any extra explanations, apologies, or "
        "introductory/concluding remarks. "
        "Only provide the ID followed by the translated text for each item.\n\n"
        "Texts to translate:\n"
    )
    return prompt_text, formatted_text_for_llm


# Helper function to apply translations
def _apply_translations_to_runs(
    all_processed_run_details: list[dict[str, Any]],
) -> None:
    """Apply the final translations back to the presentation's text runs."""
    click.echo(
        "Replacing text with translations on slides in the copied presentation..."
    )
    for item in all_processed_run_details:
        if item["final_translation"] is not None:
            item["run_object"].text = item["final_translation"]
        elif not item.get(
            "from_cache", False
        ):  # Only warn if it was supposed to be LLM translated
            click.echo(
                f"Warning: No translation found for run with original text "
                f"'{item['original_text'][:30]}...' "
                f"(LLM ID: {item.get('llm_id', 'N/A')}). Leaving original.",
                err=True,
            )


# Mode-specific processing function for "translate"
def _process_translation_mode(
    slides_to_process: list[PptxSlide],
    original_page_indices: list[int],
    cache_file_path: str,
    eol_marker: str,
    text_id_counter_start: int,
) -> int:
    """Handle the entire translation process for selected slides."""
    current_text_id_counter = text_id_counter_start
    translation_cache = load_cache(cache_file_path)

    global_texts_for_llm_prompt: list[dict[str, Any]] = []
    all_processed_run_details: list[dict[str, Any]] = []
    pending_page_cache_updates: dict[str, list[dict[str, str]]] = {}

    if slides_to_process:
        click.echo(
            f"Processing {len(slides_to_process)} selected slide(s) for translation..."
        )
        for slide_idx, slide_to_extract in enumerate(slides_to_process):
            original_page_0_indexed = original_page_indices[slide_idx]
            current_page_num_1_indexed = original_page_0_indexed + 1

            run_info_on_slide = _extract_run_info_from_slide(slide_to_extract)

            if not run_info_on_slide:
                click.echo(
                    f"  Slide {slide_idx + 1} "
                    f"(Original page {current_page_num_1_indexed}): No text found."
                )
                continue

            current_page_texts_for_hash = [
                item["original_text"] for item in run_info_on_slide
            ]
            page_hash = generate_page_hash(current_page_texts_for_hash)

            (
                texts_for_llm_slide,
                processed_runs_slide,
                current_text_id_counter,
                page_requires_llm,
            ) = prepare_slide_for_translation(
                run_info_on_slide,
                page_hash,
                translation_cache,
                current_text_id_counter,
                eol_marker,
                current_page_num_1_indexed,
            )
            global_texts_for_llm_prompt.extend(texts_for_llm_slide)
            all_processed_run_details.extend(processed_runs_slide)

            if page_requires_llm and page_hash not in pending_page_cache_updates:
                pending_page_cache_updates[page_hash] = []

    if not all_processed_run_details:
        click.echo("No text found to process on selected slides for mode 'translate'.")
        commit_pending_cache_updates(
            translation_cache, pending_page_cache_updates, cache_file_path
        )
        return current_text_id_counter

    click.echo(
        f"Processed {len(all_processed_run_details)} total text runs. "
        f"{len(global_texts_for_llm_prompt)} to translate via LLM."
    )

    if not global_texts_for_llm_prompt:
        click.echo("All text elements found in page caches. Skipping LLM prompt.")
    else:
        click.echo(
            f"Sending {len(global_texts_for_llm_prompt)} text elements to LLM for "
            "translation."
        )
        prompt_text, formatted_text_for_llm = _build_llm_prompt_and_data(
            global_texts_for_llm_prompt
        )

        click.echo("--- PROMPT TO LLM ---")
        click.echo("System/Instruction Prompt:")
        click.echo(prompt_text)
        click.echo("Data Fragments (for LLM):")
        click.echo(formatted_text_for_llm)
        click.echo("--- END OF PROMPT ---")

        model_instance = llm.get_model()
        response = model_instance.prompt(
            prompt_text, fragments=[formatted_text_for_llm]
        )
        translated_text_response_str = response.text()

        click.echo("--- RESPONSE FROM LLM ---")
        click.echo(translated_text_response_str)
        click.echo("--- END OF RESPONSE ---")
        click.echo("Received translation from LLM.")

        update_data_from_llm_response(
            translated_text_response_str.splitlines(),
            global_texts_for_llm_prompt,
            all_processed_run_details,
            pending_page_cache_updates,
            eol_marker,
        )

    _apply_translations_to_runs(all_processed_run_details)
    commit_pending_cache_updates(
        translation_cache, pending_page_cache_updates, cache_file_path
    )
    return current_text_id_counter


# Mode-specific processing function for "reverse-words"
def _process_reverse_words_mode(
    slides_to_process: list[PptxSlide],
    original_page_indices: list[int],
    eol_marker: str,
    text_id_counter_start: int,
) -> int:
    """Handle the reverse words process for selected slides."""
    current_text_id_counter = text_id_counter_start
    text_elements_for_reverse: list[dict[str, Any]] = []

    if slides_to_process:
        click.echo(
            f"Extracting text from {len(slides_to_process)} slides for "
            "mode 'reverse-words'..."
        )
        for slide_idx, slide_to_extract in enumerate(slides_to_process):
            original_page_0_indexed = original_page_indices[slide_idx]
            current_page_num_1_indexed = original_page_0_indexed + 1

            run_info_on_slide = _extract_run_info_from_slide(slide_to_extract)

            for run_detail in run_info_on_slide:
                original_text = run_detail["original_text"]
                text_with_eol = original_text + eol_marker
                text_id = f"pg{current_page_num_1_indexed}_txt{current_text_id_counter}"
                text_elements_for_reverse.append(
                    {
                        "id": text_id,
                        "text_with_eol_for_reverse": text_with_eol,
                        "run_object": run_detail["run_object"],
                    }
                )
                current_text_id_counter += 1

    if not text_elements_for_reverse:
        click.echo(
            "No text found to process on selected slides for mode 'reverse-words'."
        )
    else:
        click.echo(
            f"Found {len(text_elements_for_reverse)} text elements to process for "
            "mode 'reverse-words'."
        )
        click.echo("Applying word reversal on slides...")
        for item in text_elements_for_reverse:
            reversed_text_with_eol = reverse_individual_words(
                item["text_with_eol_for_reverse"]
            )
            final_reversed_text = reversed_text_with_eol.removesuffix(eol_marker)
            item["run_object"].text = final_reversed_text
        click.echo("Text replaced with reversed-word text on slides.")
    return current_text_id_counter


@click.command()
@click.option(
    "--mode",
    type=click.Choice(
        ["translate", "reverse-words"],
        case_sensitive=False,
    ),
    default="translate",
    show_default=True,
    help="Operation mode for the script.",
)
@click.option(
    "--pages",
    type=str,
    default=None,
    help=(
        "Specify page range to process (e.g., '1,3-5,8-'). 1-indexed. "
        "Processes all pages if not specified."
    ),
)
@click.argument("input_path", type=click.Path(exists=True, dir_okay=False))
@click.argument("output_path", type=click.Path(dir_okay=False))
def main(input_path: str, output_path: str, mode: str, pages: str | None) -> None:
    """Process a PowerPoint presentation.

    It first copies the input presentation to the output path.
    Then, for 'translate' and 'reverse-words' modes, text on selected slides
    within this copied presentation is modified in place.
    'translate' mode uses an on-disk cache.
    Page range can be specified using the --pages option.
    """
    cache_file_path = "translation_cache.json"

    click.echo(f"Copying '{input_path}' to '{output_path}' to preserve layout...")
    shutil.copy2(input_path, output_path)
    click.echo("File copy complete.")

    click.echo(f"Loading presentation for modification from: {output_path}")
    prs = Presentation(output_path)

    num_original_slides = len(prs.slides)
    if num_original_slides == 0:
        click.echo("Input presentation has no slides. Exiting.")
        return

    selected_pages_0_indexed = _handle_slide_selection(pages, num_original_slides)

    slides_to_process_objects: list[PptxSlide] = []
    original_page_indices_for_processing: list[int] = []

    for i, slide_obj in enumerate(prs.slides):
        if i in selected_pages_0_indexed:
            slides_to_process_objects.append(slide_obj)
            original_page_indices_for_processing.append(i)

    if not slides_to_process_objects:
        if pages:
            click.echo(
                f"Warning: The specified page range '{pages}' resulted in "
                f"no slides being selected from {num_original_slides} total slides. "
                "No text processing will occur."
            )
        else:
            click.echo(
                f"No slides selected for processing from {num_original_slides} "
                "total slides. No text processing will occur."
            )

        prs.save(output_path)  # Save presentation even if no slides processed
        click.echo(f"Presentation saved (no text modifications) to: {output_path}")
        if mode == "translate":
            translation_cache = load_cache(cache_file_path)
            commit_pending_cache_updates(translation_cache, {}, cache_file_path)
        return

    click.echo(
        f"Preparing to process text on {len(slides_to_process_objects)} selected "
        f"slides (out of {num_original_slides} total) in the copied presentation..."
    )

    text_id_counter = 0

    if mode == "translate":
        text_id_counter = _process_translation_mode(
            slides_to_process_objects,
            original_page_indices_for_processing,
            cache_file_path,
            EOL_MARKER,
            text_id_counter,
        )
    elif mode == "reverse-words":
        text_id_counter = _process_reverse_words_mode(
            slides_to_process_objects,
            original_page_indices_for_processing,
            EOL_MARKER,
            text_id_counter,
        )
    else:  # pragma: no cover
        # This case should be prevented by click.Choice on the --mode option.
        # If execution reaches here, it's an unexpected state.
        # For robustness, one might raise an error, but relying on Click is typical.
        pass

    prs.save(output_path)
    _emit_save_message(mode, Path(output_path))


if __name__ == "__main__":
    main()
