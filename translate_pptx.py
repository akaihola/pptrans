import copy  # For deepcopying elements
import json
import os
import shutil  # For file copying

import click
import llm  # Simon Willison's LLM library
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.util import Inches, Pt

EOL_MARKER = "<"
# import xml.etree.ElementTree as ET # Not strictly needed if using copy.deepcopy for lxml


def load_cache(cache_file_path):
    """Loads the translation cache from a JSON file."""
    if os.path.exists(cache_file_path):
        try:
            with open(cache_file_path, "r", encoding="utf-8") as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError) as e:
            click.echo(
                f"Warning: Could not load cache file {cache_file_path}. Error: {e}. Starting with an empty cache.",
                err=True,
            )
    return {}


def save_cache(cache_data, cache_file_path):
    """Saves the translation cache to a JSON file."""
    try:
        with open(cache_file_path, "w", encoding="utf-8") as f:
            json.dump(cache_data, f, indent=4, ensure_ascii=False)
    except IOError as e:
        click.echo(
            f"Warning: Could not save cache file {cache_file_path}. Error: {e}",
            err=True,
        )


def reverse_individual_words(text_string_with_eol):
    """
    Reverses each word in a space-separated string, preserving an EOL_MARKER if present.
    """
    text_to_reverse = text_string_with_eol
    had_eol = False
    if text_string_with_eol.endswith(EOL_MARKER):
        text_to_reverse = text_string_with_eol[: -len(EOL_MARKER)]
        had_eol = True

    words = text_to_reverse.split(" ")
    reversed_words = [word[::-1] for word in words]
    result = " ".join(reversed_words)

    if had_eol:
        return result + EOL_MARKER
    return result


def parse_page_range(range_str, total_slides):
    """
    Parses a page range string (e.g., "1,3-5,8-") into a set of 0-indexed page numbers.
    Validates against total_slides. Raises click.BadParameter on error.
    1-indexed page numbers are used in the input string.
    """
    selected_pages_0_indexed = set()

    if total_slides == 0:
        # If a range_str was provided for a presentation with no slides, it's an issue.
        # If range_str is None or empty, this function might not even be called,
        # or it could return an empty set, which is fine.
        if range_str:  # Check if user actually specified a range for an empty deck
            raise click.BadParameter(
                f"Cannot specify page range '{range_str}' for a presentation with no slides."
            )
        return selected_pages_0_indexed  # Correctly returns empty set if no slides

    # If range_str is None (meaning --pages not used), this function shouldn't be called by main.
    # The caller (main) should handle that by selecting all pages.
    # If range_str is an empty string (e.g., --pages=""), it's an invalid input if --pages was explicitly used.
    if (
        range_str is not None and not range_str.strip()
    ):  # Check if --pages was used and an empty string was passed
        raise click.BadParameter(
            "Page range string cannot be empty if --pages option is used with an empty value."
        )

    parts = range_str.split(",")
    for part_specifier in parts:
        part = part_specifier.strip()
        if not part:  # Handles cases like "1,,2" or leading/trailing commas
            continue

        if "-" in part:
            # Range specifier
            elements = part.split("-", 1)
            start_str, end_str = elements[0].strip(), elements[1].strip()

            start_page_1_indexed = None
            end_page_1_indexed = None

            if start_str:  # "N-M" or "N-"
                try:
                    start_page_1_indexed = int(start_str)
                    if start_page_1_indexed < 1:
                        raise click.BadParameter(
                            f"Page numbers must be positive. Found start '{start_str}' in '{part_specifier}'."
                        )
                except ValueError:
                    raise click.BadParameter(
                        f"Invalid start page number '{start_str}' in '{part_specifier}'. Must be an integer."
                    )
            else:  # "-M"
                start_page_1_indexed = 1  # Default start for ranges like "-5"

            if end_str:  # "N-M" or "-M"
                try:
                    end_page_1_indexed = int(end_str)
                    if end_page_1_indexed < 1:
                        raise click.BadParameter(
                            f"Page numbers must be positive. Found end '{end_str}' in '{part_specifier}'."
                        )
                except ValueError:
                    raise click.BadParameter(
                        f"Invalid end page number '{end_str}' in '{part_specifier}'. Must be an integer."
                    )
            else:  # "N-"
                end_page_1_indexed = total_slides  # Default end for ranges like "5-"

            if start_page_1_indexed > end_page_1_indexed:
                raise click.BadParameter(
                    f"Start page {start_page_1_indexed} cannot be greater than end page {end_page_1_indexed} in range '{part_specifier}'."
                )

            # Validate individual bounds against total_slides
            if start_page_1_indexed > total_slides:
                # This also covers cases where start_page_1_indexed was 1 (for "-M") but total_slides is 0 (already handled)
                # or if user specifies "10-" for a 5-slide deck.
                # No need to add to set if start is already beyond total slides.
                click.echo(
                    f"Warning: Start page {start_page_1_indexed} in range '{part_specifier}' is beyond the total of {total_slides} slides. This part of the range will select no pages.",
                    err=True,
                )
                continue  # Skip this part of the range

            # For "N-M" or "-M", if end_page_1_indexed was explicitly given and exceeds total_slides
            if end_str and end_page_1_indexed > total_slides:
                click.echo(
                    f"Warning: End page {end_page_1_indexed} in range '{part_specifier}' is beyond the total of {total_slides} slides. Range will be capped at {total_slides}.",
                    err=True,
                )
                # end_page_1_indexed will be capped by min() below.

            # Add pages (0-indexed) to the set
            # Cap end_page_1_indexed at total_slides for ranges like "N-" or valid "N-M"
            # This loop will correctly handle start_page_1_indexed up to min(end_page_1_indexed, total_slides)
            for i in range(
                start_page_1_indexed, min(end_page_1_indexed, total_slides) + 1
            ):
                selected_pages_0_indexed.add(i - 1)
        else:
            # Single page specifier
            try:
                page_num_1_indexed = int(part)
                if not (1 <= page_num_1_indexed <= total_slides):
                    raise click.BadParameter(
                        f"Page number {page_num_1_indexed} in '{part_specifier}' is out of valid range [1, {total_slides}]."
                    )
                selected_pages_0_indexed.add(page_num_1_indexed - 1)
            except ValueError:
                raise click.BadParameter(
                    f"Invalid page number: '{part_specifier}'. Must be an integer."
                )

    return selected_pages_0_indexed


@click.command()
@click.option(
    "--mode",
    type=click.Choice(
        ["translate", "reverse-words"],
        case_sensitive=False,  # 'duplicate-only' mode removed
    ),
    default="translate",
    show_default=True,
    help="Operation mode for the script.",
)
@click.option(
    "--pages",
    type=str,
    default=None,
    help="Specify page range to process (e.g., '1,3-5,8-'). 1-indexed. Processes all pages if not specified.",
)
@click.argument("input_path", type=click.Path(exists=True, dir_okay=False))
@click.argument("output_path", type=click.Path(dir_okay=False))
def main(input_path, output_path, mode, pages):
    """
    Processes a PowerPoint presentation.
    It first copies the input presentation to the output path.
    Then, for 'translate' and 'reverse-words' modes, text on selected slides
    within this copied presentation is modified in place.
    'translate' mode uses an on-disk cache.
    Page range can be specified using the --pages option.
    """
    cache_file_path = (
        "translation_cache.json"  # Cache file in the same directory as script execution
    )

    click.echo(f"Copying '{input_path}' to '{output_path}' to preserve layout...")
    try:
        shutil.copy2(input_path, output_path)
    except Exception as e:
        click.echo(f"Error copying file: {e}", err=True)
        return
    click.echo("File copy complete.")

    click.echo(f"Loading presentation for modification from: {output_path}")
    prs = Presentation(output_path)

    num_original_slides = len(prs.slides)
    if num_original_slides == 0:
        click.echo("Input presentation has no slides. Exiting.")
        return

    selected_pages_0_indexed = set()
    if pages:
        try:
            selected_pages_0_indexed = parse_page_range(pages, num_original_slides)
        except click.BadParameter as e:
            # parse_page_range raises BadParameter, which click handles by exiting.
            # We can re-raise if we want to be explicit or add more context, but click does it.
            # For now, let click handle the exit.
            # click.echo(f"Error: {e}", err=True) # This would be redundant
            raise  # Re-raise to ensure click handles it as expected
    else:
        # If --pages is not specified, process all slides
        selected_pages_0_indexed = set(range(num_original_slides))

    slides_to_process_objects = []
    if num_original_slides > 0:
        for i, slide_obj in enumerate(prs.slides):
            if i in selected_pages_0_indexed:
                slides_to_process_objects.append(slide_obj)

    if not slides_to_process_objects:
        if pages:  # User specified a range, but it resulted in no slides
            click.echo(
                f"Warning: The specified page range '{pages}' resulted in no slides being selected from {num_original_slides} total slides. No text processing will occur."
            )
        else:  # No pages specified, but presentation was empty (already handled) or became empty (should not happen here)
            click.echo(
                f"No slides selected for processing from {num_original_slides} total slides. No text processing will occur."
            )
        # Save the copied presentation and exit if no slides are to be processed.
        # This is important if only a copy was intended or if the range was valid but empty.
        prs.save(output_path)
        click.echo(f"Presentation saved (no text modifications) to: {output_path}")
        # Save cache even if no text elements were processed, in case it was loaded and needs to be preserved.
        if mode == "translate":
            translation_cache = load_cache(
                cache_file_path
            )  # Ensure cache is loaded if not already
            save_cache(translation_cache, cache_file_path)
        return

    # Renaming slides_for_text_extraction to slides_to_process_objects for clarity
    # The old slides_for_text_extraction is now slides_to_process_objects

    # The original logic for 'slides_for_text_extraction' initialization is now replaced by the above.
    # The following 'if mode == "translate" or mode == "reverse-words":' block's content
    # related to populating slides_for_text_extraction is no longer needed here as slides_to_process_objects is already set.
    # We just need to update the logging.

    if mode == "translate" or mode == "reverse-words":
        # slides_for_text_extraction = list(prs.slides) # This line is now replaced by slides_to_process_objects logic
        if slides_to_process_objects:  # Only print if there are slides to process
            click.echo(
                f"Preparing to process text on {len(slides_to_process_objects)} selected slide(s) (out of {num_original_slides} total) in the copied presentation..."
            )
    # 'duplicate-only' mode and slide duplication logic removed.
    # Text processing will now occur directly on the slides in 'slides_for_text_extraction'.

    # Text processing for 'translate' and 'reverse-words' modes
    text_id_counter = 0  # Shared counter for text element IDs

    if mode == "translate":
        click.echo(f"Loading translation cache from: {cache_file_path}")
        translation_cache = load_cache(cache_file_path)

        texts_for_llm_prompt = []
        all_text_elements_with_status = []

        if slides_to_process_objects:
            click.echo(
                f"Extracting text from {len(slides_to_process_objects)} slides in the copied presentation for mode '{mode}' (with cache checking)..."
            )
            for slide_to_extract in slides_to_process_objects:
                for shape in slide_to_extract.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text  # No strip. This is the canonical version for cache key and internal "text" storage.
                                if original_text:  # Check unstripped text
                                    text_id = f"text_{text_id_counter}"
                                    text_id_counter += 1

                                    # Cache uses original_text (no EOL) as key
                                    if original_text in translation_cache:
                                        click.echo(
                                            f"  Cache hit for ID {text_id}: '{original_text[:30].replace(chr(10), ' ').replace(chr(13), ' ')}...'"
                                        )
                                        all_text_elements_with_status.append(
                                            {
                                                "id": text_id,
                                                "text": original_text,  # Store raw original_text
                                                "run_object": run,
                                                "translation": translation_cache[
                                                    original_text
                                                ],  # Translation from cache is EOL-stripped
                                                "from_cache": True,
                                            }
                                        )
                                    else:  # Cache miss
                                        click.echo(
                                            f"  Cache miss for ID {text_id}: '{original_text[:30].replace(chr(10), ' ').replace(chr(13), ' ')}...' (will send to LLM)"
                                        )
                                        text_to_send_to_llm = original_text + EOL_MARKER
                                        texts_for_llm_prompt.append(
                                            {
                                                "id": text_id,
                                                "original_text_for_cache": original_text,  # Key for cache update later
                                                "text_to_send": text_to_send_to_llm,  # Actual text for LLM
                                            }
                                        )
                                        all_text_elements_with_status.append(
                                            {
                                                "id": text_id,
                                                "text": original_text,  # Store raw original_text
                                                "run_object": run,
                                                "translation": None,  # Will be filled (EOL-stripped)
                                                "from_cache": False,
                                            }
                                        )
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        original_text = (
                                            run.text
                                        )  # No strip. Canonical for cache key.
                                        if original_text:
                                            text_id = f"text_{text_id_counter}"
                                            text_id_counter += 1
                                            # Cache uses original_text (no EOL) as key
                                            if original_text in translation_cache:
                                                click.echo(
                                                    f"  Cache hit for ID {text_id} (table): '{original_text[:30].replace(chr(10), ' ').replace(chr(13), ' ')}...'"
                                                )
                                                all_text_elements_with_status.append(
                                                    {
                                                        "id": text_id,
                                                        "text": original_text,  # Store raw original_text
                                                        "run_object": run,
                                                        "translation": translation_cache[
                                                            original_text
                                                        ],  # EOL-stripped from cache
                                                        "from_cache": True,
                                                    }
                                                )
                                            else:  # Cache miss
                                                click.echo(
                                                    f"  Cache miss for ID {text_id} (table): '{original_text[:30].replace(chr(10), ' ').replace(chr(13), ' ')}...' (will send to LLM)"
                                                )
                                                text_to_send_to_llm = (
                                                    original_text + EOL_MARKER
                                                )
                                                texts_for_llm_prompt.append(
                                                    {
                                                        "id": text_id,
                                                        "original_text_for_cache": original_text,  # Key for cache update
                                                        "text_to_send": text_to_send_to_llm,  # Actual text for LLM
                                                    }
                                                )
                                                all_text_elements_with_status.append(
                                                    {
                                                        "id": text_id,
                                                        "text": original_text,  # Store raw original_text
                                                        "run_object": run,
                                                        "translation": None,  # Will be EOL-stripped
                                                        "from_cache": False,
                                                    }
                                                )

        if not all_text_elements_with_status:
            click.echo(
                f"No text found to process on slides in the copied presentation for mode '{mode}'."
            )
            prs.save(output_path)
            click.echo(
                f"Presentation saved without text modification in '{mode}' mode to: {output_path}"
            )
            click.echo(f"Saving cache (even if empty/unchanged) to: {cache_file_path}")
            save_cache(
                translation_cache, cache_file_path
            )  # Save cache even if no text elements
            return

        click.echo(
            f"Found {len(all_text_elements_with_status)} total text elements. {len(texts_for_llm_prompt)} to translate via LLM."
        )

        if not texts_for_llm_prompt:
            click.echo("All text elements found in cache. Skipping LLM prompt.")
        else:
            click.echo(
                f"Sending {len(texts_for_llm_prompt)} text elements to LLM for translation."
            )
            formatted_text_for_llm = "\n".join(
                [
                    f"{item['id']}:{item['text_to_send']}"
                    for item in texts_for_llm_prompt
                ]  # Use text_to_send
            )
            prompt_text = (
                "You are an expert Finnish to English translator. "
                "Translate the following text segments accurately from Finnish to English. "
                "Each segment is prefixed with a unique ID (e.g., text_0, text_1). "
                "IMPORTANT: A sequence of text items (e.g., text_0, text_1, text_2) may represent a single continuous sentence that has been split due to formatting. Interpret and translate such sequences as a coherent whole sentence to maintain context and flow. "
                "The text for each ID might end with an EOL marker: '<'. "  # Added EOL marker info
                "Your response MUST consist ONLY of the translated segments, each prefixed with its original ID, "
                "and each on a new line. Maintain the exact ID and format. "
                "PRESERVE ALL LEADING AND TRAILING WHITESPACE from the original segment in your translation. "  # Added whitespace preservation
                "If an EOL marker '<' was present at the end of the input segment, IT MUST be present at the end of your translated segment, including any whitespace before it.\n"  # Added EOL marker preservation
                "For example, if you receive:\n"
                "text_0: Tämä on pitkä \n"
                "text_1:lause, joka on \n"
                "text_2:jaettu.<\n"
                "text_3:    Toinen lause.   <\n"
                "text_4: Yksittäinen.\n"
                "You MUST return:\n"
                "text_0: This is a long \n"
                "text_1:sentence that has been \n"
                "text_2:split.<\n"
                "text_3:    Another sentence.   <\n"
                "text_4: Standalone.\n\n"
                "Do not add any extra explanations, apologies, or introductory/concluding remarks. "
                "Only provide the ID followed by the translated text for each item.\n\n"
                "Texts to translate:\n"
            )
            click.echo("--- PROMPT TO LLM ---")
            click.echo("System/Instruction Prompt:")
            click.echo(prompt_text)
            click.echo("Data Fragments (for LLM):")
            click.echo(formatted_text_for_llm)
            click.echo("--- END OF PROMPT ---")

            model_instance = llm.get_model()
            response = model_instance.prompt(
                prompt_text, fragments=[formatted_text_for_llm]
            )
            translated_text_response = response.text()

            click.echo("--- RESPONSE FROM LLM ---")
            click.echo(translated_text_response)
            click.echo("--- END OF RESPONSE ---")
            click.echo("Received translation from LLM.")

            for line in translated_text_response.splitlines():
                line = (
                    line.strip()
                )  # Strip the whole line to check if it's empty or just whitespace
                if not line:
                    continue
                try:
                    parts = line.split(":", 1)
                    if len(parts) == 2:
                        parsed_text_id = parts[0].strip()
                        # DO NOT strip llm_translation here, preserve spaces from LLM
                        llm_translation_with_eol = parts[1]

                        # Find the corresponding original_text_for_cache from texts_for_llm_prompt
                        # This original_text_for_cache is what we use as the key in translation_cache
                        prompt_item_data = next(
                            (
                                item
                                for item in texts_for_llm_prompt
                                if item["id"] == parsed_text_id
                            ),
                            None,
                        )

                        if prompt_item_data:
                            original_text_for_cache_key = prompt_item_data[
                                "original_text_for_cache"
                            ]

                            # Strip EOL_MARKER from the translation before caching and storing
                            final_llm_translation = llm_translation_with_eol
                            if final_llm_translation.endswith(EOL_MARKER):
                                final_llm_translation = final_llm_translation[
                                    : -len(EOL_MARKER)
                                ]

                            translation_cache[original_text_for_cache_key] = (
                                final_llm_translation
                            )

                            for elem in all_text_elements_with_status:
                                if elem["id"] == parsed_text_id:
                                    elem["translation"] = (
                                        final_llm_translation  # Store EOL-stripped
                                    )
                                    elem["from_cache"] = False
                                    break
                        else:
                            click.echo(
                                f"Warning: Could not find original_text_for_cache for ID {parsed_text_id} from LLM response to update cache/status list.",
                                err=True,
                            )
                    else:
                        click.echo(
                            f"Warning: Could not parse translation line: {line}",
                            err=True,
                        )
                except Exception as e:
                    click.echo(
                        f"Warning: Error parsing translation line '{line}': {e}",
                        err=True,
                    )

        click.echo(
            "Replacing text with translations on slides in the copied presentation..."
        )
        for item in all_text_elements_with_status:
            if item["translation"]:
                item["run_object"].text = item["translation"]

        click.echo(f"Saving updated translation cache to: {cache_file_path}")
        save_cache(translation_cache, cache_file_path)

    elif mode == "reverse-words":
        # Original logic for reverse-words, uses a simple text_elements list
        text_elements_for_reverse = []
        if slides_to_process_objects:
            click.echo(
                f"Extracting text from {len(slides_to_process_objects)} slides in the copied presentation for mode '{mode}'..."
            )
            for slide_to_extract in slides_to_process_objects:
                for shape in slide_to_extract.shapes:
                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            for run in paragraph.runs:
                                original_text = run.text  # No strip
                                if original_text:
                                    text_with_eol = original_text + EOL_MARKER
                                    text_id = (
                                        f"text_{text_id_counter}"  # Uses global counter
                                    )
                                    text_elements_for_reverse.append(
                                        {
                                            "id": text_id,
                                            "text": text_with_eol,  # Store with EOL
                                            "run_object": run,
                                        }
                                    )
                                    text_id_counter += 1
                    if shape.has_table:
                        for row in shape.table.rows:
                            for cell in row.cells:
                                for paragraph in cell.text_frame.paragraphs:
                                    for run in paragraph.runs:
                                        original_text = run.text  # No strip
                                        if original_text:
                                            text_with_eol = original_text + EOL_MARKER
                                            text_id = f"text_{text_id_counter}"  # Uses global counter
                                            text_elements_for_reverse.append(
                                                {
                                                    "id": text_id,
                                                    "text": text_with_eol,  # Store with EOL
                                                    "run_object": run,
                                                }
                                            )
                                            text_id_counter += 1

        if not text_elements_for_reverse:
            click.echo(
                f"No text found to process on slides in the copied presentation for mode '{mode}'."
            )
            prs.save(output_path)
            click.echo(
                f"Presentation saved without text modification in '{mode}' mode to: {output_path}"
            )
            return

        click.echo(
            f"Found {len(text_elements_for_reverse)} text elements to process for mode '{mode}'."
        )
        click.echo("Applying word reversal on slides in the copied presentation...")
        for item in text_elements_for_reverse:
            reversed_text_with_eol = reverse_individual_words(item["text"])
            final_reversed_text = reversed_text_with_eol
            if final_reversed_text.endswith(EOL_MARKER):
                final_reversed_text = final_reversed_text[: -len(EOL_MARKER)]
            item["run_object"].text = final_reversed_text
        click.echo(
            "Text replaced with reversed-word text on slides in the copied presentation."
        )

    prs.save(output_path)
    click.echo(f"Presentation saved in '{mode}' mode to: {output_path}")


if __name__ == "__main__":
    main()
