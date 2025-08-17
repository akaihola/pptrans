# /// script
# requires-python = ">=3.11"
# dependencies = [
#     "python-pptx==1.0.2",
# ]
# ///

import marimo

__generated_with = "0.14.17"
app = marimo.App(width="medium")


@app.cell
def _():
    from pathlib import Path

    import marimo as mo
    import pptx
    from pptx.util import Inches
    return Inches, Path, pptx


@app.cell
def _(Inches, Path, pptx):
    prs = pptx.Presentation()
    blank_slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_slide_layout)

    left = Inches(1.0)
    top = Inches(1.0)
    width = Inches(8.0)
    height = Inches(1.0)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "Translation from English to Finnish is working perfectly."

    prs.save(Path(__file__).parent / "sample.en.pptx")
    return


@app.cell
def _():
    return


if __name__ == "__main__":
    app.run()
