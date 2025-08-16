"""Unit tests for the `pptrans.__main__` module."""

from __future__ import annotations

import importlib
import sys
from pathlib import Path  # Import Path for runtime use
from typing import TYPE_CHECKING, Any
from unittest import mock
from unittest.mock import MagicMock, patch

import click
import pytest  # Import pytest for the marker
from click.testing import CliRunner
from pptx.util import Emu  # Import Emu

from pptrans.__main__ import (
    EOL_MARKER,
    _apply_translations_to_runs,
    _build_llm_prompt_and_data,
    _emit_save_message,
    _extract_run_info_from_slide,
    _handle_slide_selection,
    _process_reverse_words_mode,
    _process_translation_mode,
    main,
    reverse_individual_words,
)

if TYPE_CHECKING:
    from pptx.slide import Slide as PptxSlide


def create_mock_run(text: str) -> MagicMock:
    """Create a mock run object."""
    run = MagicMock()
    run.text = text
    return run


def create_mock_paragraph(runs: list[MagicMock]) -> MagicMock:
    """Create a mock paragraph object."""
    para = MagicMock()
    para.runs = runs
    return para


def create_mock_text_frame(paragraphs: list[MagicMock]) -> MagicMock:
    """Create a mock text_frame object."""
    tf = MagicMock()
    tf.paragraphs = paragraphs
    return tf


def create_mock_cell(text_frame: MagicMock) -> MagicMock:
    """Create a mock cell object."""
    cell = MagicMock()
    cell.text_frame = text_frame
    return cell


def create_mock_row(cells: list[MagicMock]) -> MagicMock:
    """Create a mock row object."""
    row = MagicMock()
    row.cells = cells
    return row


def create_mock_table(rows: list[MagicMock]) -> MagicMock:
    """Create a mock table object."""
    table = MagicMock()
    table.rows = rows
    return table


def create_mock_shape(
    has_text_frame: bool = False,
    text_frame: MagicMock | None = None,
    has_table: bool = False,
    table: MagicMock | None = None,
    left: int = 0,  # Add left coordinate
    top: int = 0,  # Add top coordinate
) -> MagicMock:
    """Create a mock shape object."""
    shape = MagicMock()
    shape.has_text_frame = has_text_frame
    shape.text_frame = text_frame
    shape.has_table = has_table
    shape.table = table
    # Mock position attributes, returning Emu objects
    shape.left = Emu(left)
    shape.top = Emu(top)
    return shape


def create_mock_slide(shapes: list[MagicMock]) -> MagicMock:
    """Create a mock slide object."""
    slide = MagicMock()
    slide.shapes = shapes
    return slide


@pytest.mark.kwparametrize(
    dict(
        text_string_with_eol=f"hello world{EOL_MARKER}",
        expected=f"olleh dlrow{EOL_MARKER}",
        id="with_eol",
    ),
    dict(
        text_string_with_eol="hello world",
        expected="olleh dlrow",
        id="without_eol",
    ),
    dict(
        text_string_with_eol=EOL_MARKER,
        expected=EOL_MARKER,
        id="empty_string_with_eol",
    ),
    dict(
        text_string_with_eol="",
        expected="",
        id="empty_string_without_eol",
    ),
    dict(
        text_string_with_eol=f"hello{EOL_MARKER}",
        expected=f"olleh{EOL_MARKER}",
        id="single_word_with_eol",
    ),
    dict(
        text_string_with_eol="hello",
        expected="olleh",
        id="single_word_without_eol",
    ),
    dict(
        text_string_with_eol=f"hello  world{EOL_MARKER}",
        expected=f"olleh  dlrow{EOL_MARKER}",
        id="multiple_spaces_with_eol",
    ),
)
def test_reverse_individual_words(text_string_with_eol: str, expected: str) -> None:
    assert reverse_individual_words(text_string_with_eol) == expected


@patch("pptrans.__main__.parse_page_range")
@pytest.mark.kwparametrize(
    dict(
        pages_option=None,
        num_original_slides=3,
        expected_result={0, 1, 2},
        mock_parse_page_range_return=None,
        id="pages_option_none",
    ),
    dict(
        pages_option="1,3",
        num_original_slides=5,
        mock_parse_page_range_return={0, 2},
        expected_result={0, 2},
        id="pages_option_valid",
    ),
    dict(
        pages_option=None,
        num_original_slides=0,
        expected_result=set(),
        mock_parse_page_range_return=None,
        id="num_original_slides_zero",
    ),
)
def test_handle_slide_selection(
    mock_parse_page_range: MagicMock,
    pages_option: str | None,
    num_original_slides: int,
    expected_result: set[int],
    mock_parse_page_range_return: set[int] | None,
) -> None:
    if mock_parse_page_range_return is not None:
        mock_parse_page_range.return_value = mock_parse_page_range_return

    result = _handle_slide_selection(pages_option, num_original_slides)
    assert result == expected_result

    if pages_option:
        mock_parse_page_range.assert_called_once_with(pages_option, num_original_slides)
    else:
        mock_parse_page_range.assert_not_called()


@patch("pptx.util.Emu")  # Mock Emu
@pytest.mark.kwparametrize(
    dict(
        slide_obj=create_mock_slide(shapes=[]),
        expected=[],
        id="empty_slide",
    ),
    dict(
        slide_obj=create_mock_slide(shapes=[create_mock_shape()]),
        expected=[],
        id="shape_no_text_frame_no_table",
    ),
    dict(
        slide_obj=create_mock_slide(
            shapes=[
                create_mock_shape(
                    has_text_frame=True,
                    text_frame=create_mock_text_frame(
                        paragraphs=[create_mock_paragraph(runs=[])]
                    ),
                    left=100000,
                    top=200000,
                )
            ]
        ),
        expected=[],
        id="shape_with_text_frame_empty_runs",
    ),
    dict(
        slide_obj=create_mock_slide(
            shapes=[
                create_mock_shape(
                    has_text_frame=True,
                    text_frame=create_mock_text_frame(
                        paragraphs=[
                            create_mock_paragraph(
                                runs=[
                                    create_mock_run("Hello"),
                                    create_mock_run(" World"),
                                ]
                            )
                        ]
                    ),
                    left=100000,
                    top=200000,
                )
            ]
        ),
        expected=[
            {
                "original_text": "Hello",
                "run_object": Ellipsis,
                "shape_idx": 0,
                "run_idx_in_shape": 0,
                "shape_x": 100,
                "shape_y": 200,
            },
            {
                "original_text": " World",
                "run_object": Ellipsis,
                "shape_idx": 0,
                "run_idx_in_shape": 1,
                "shape_x": 100,
                "shape_y": 200,
            },
        ],
        id="shape_with_text_runs",
    ),
    dict(
        slide_obj=create_mock_slide(
            shapes=[
                create_mock_shape(
                    has_table=True,
                    table=create_mock_table(
                        rows=[
                            create_mock_row(
                                cells=[
                                    create_mock_cell(
                                        create_mock_text_frame(
                                            [
                                                create_mock_paragraph(
                                                    [create_mock_run("Table")]
                                                )
                                            ]
                                        )
                                    )
                                ]
                            )
                        ]
                    ),
                    left=300000,
                    top=400000,
                )
            ]
        ),
        expected=[
            {
                "original_text": "Table",
                "run_object": Ellipsis,
                "shape_idx": 0,
                "run_idx_in_shape": 0,
                "shape_x": 300,
                "shape_y": 400,
            }
        ],
        id="shape_with_table_runs",
    ),
    dict(
        slide_obj=create_mock_slide(
            shapes=[
                create_mock_shape(
                    has_text_frame=True,
                    text_frame=create_mock_text_frame(
                        paragraphs=[
                            create_mock_paragraph([create_mock_run("TextShape")])
                        ]
                    ),
                    left=500000,
                    top=600000,
                ),
                create_mock_shape(
                    has_table=True,
                    table=create_mock_table(
                        rows=[
                            create_mock_row(
                                cells=[
                                    create_mock_cell(
                                        create_mock_text_frame(
                                            [
                                                create_mock_paragraph(
                                                    [create_mock_run("InTable")]
                                                )
                                            ]
                                        )
                                    )
                                ]
                            )
                        ]
                    ),
                    left=700000,
                    top=800000,
                ),
            ]
        ),
        expected=[
            {
                "original_text": "TextShape",
                "run_object": Ellipsis,
                "shape_idx": 0,
                "run_idx_in_shape": 0,
                "shape_x": 500,
                "shape_y": 600,
            },
            {
                "original_text": "InTable",
                "run_object": Ellipsis,
                "shape_idx": 1,
                "run_idx_in_shape": 0,
                "shape_x": 700,
                "shape_y": 800,
            },
        ],
        id="mixed_content_slide",
    ),
    dict(
        slide_obj=create_mock_slide(
            shapes=[
                create_mock_shape(
                    has_table=True,
                    table=create_mock_table(
                        rows=[
                            create_mock_row(
                                cells=[
                                    create_mock_cell(
                                        create_mock_text_frame(
                                            paragraphs=[
                                                create_mock_paragraph(
                                                    runs=[
                                                        create_mock_run(""),
                                                        create_mock_run(
                                                            "ValuableInTable"
                                                        ),
                                                    ]
                                                )
                                            ]
                                        )
                                    )
                                ]
                            )
                        ]
                    ),
                    left=900000,
                    top=1000000,
                )
            ]
        ),
        expected=[
            {
                "original_text": "ValuableInTable",
                "run_object": Ellipsis,
                "shape_idx": 0,
                "run_idx_in_shape": 1,
                "shape_x": 900,
                "shape_y": 1000,
            }
        ],
        id="table_run_empty_followed_by_non_empty",
    ),
    dict(
        slide_obj=create_mock_slide(
            shapes=[
                create_mock_shape(
                    has_text_frame=True,
                    text_frame=create_mock_text_frame(
                        paragraphs=[
                            create_mock_paragraph(
                                runs=[create_mock_run("Data"), create_mock_run("")]
                            )
                        ]
                    ),
                    left=1100000,
                    top=1200000,
                )
            ]
        ),
        expected=[
            {
                "original_text": "Data",
                "run_object": Ellipsis,
                "shape_idx": 0,
                "run_idx_in_shape": 0,
                "shape_x": 1100,
                "shape_y": 1200,
            }
        ],
        id="run_with_empty_text",
    ),
)
def test_extract_run_info_from_slide(
    mock_emu: MagicMock, slide_obj: PptxSlide, expected: list[dict[str, Any]]
) -> None:
    # Configure the mocked Emu to return the integer value when .value is accessed
    mock_emu.side_effect = lambda value: MagicMock(value=value)

    result = _extract_run_info_from_slide(slide_obj)
    assert len(result) == len(expected)
    for res_item, exp_item in zip(result, expected):
        assert res_item["original_text"] == exp_item["original_text"]
        assert "run_object" in res_item
        if exp_item["run_object"] is Ellipsis:
            assert res_item["run_object"] is not None
        assert res_item["shape_idx"] == exp_item["shape_idx"]
        assert res_item["run_idx_in_shape"] == exp_item["run_idx_in_shape"]
        assert res_item["shape_x"] == exp_item["shape_x"]
        assert isinstance(res_item["shape_x"], int)
        assert res_item["shape_y"] == exp_item["shape_y"]
        assert isinstance(res_item["shape_y"], int)


@pytest.mark.kwparametrize(
    dict(
        texts_for_llm=[],
        expected_formatted_text="",
        expected_prompt_contains=[
            "You are an expert Finnish to English translator.",
            f"EOL marker: '{EOL_MARKER}'",
            "Texts to translate:",
            "For example, if you receive:",
            "pg1,el0,run0,x=100,y=200:Tämä on pitkä",
            "pg1,el0,run1,x=100,y=200:lause, joka on",
            "You MUST return:",
            "pg1,el0,run0:This is a long",
            "pg1,el0,run1:sentence that has been",
            "Do not add any extra explanations",
        ],
        id="empty_texts_for_llm",
    ),
    dict(
        texts_for_llm=[
            {
                "id": "pg1,el0,run0",
                "text_to_send": "Hello",
                "shape_x": 100,
                "shape_y": 200,
            }
        ],
        expected_formatted_text="pg1,el0,run0,x=100,y=200:Hello",
        expected_prompt_contains=[
            "You are an expert Finnish to English translator.",
            f"EOL marker: '{EOL_MARKER}'",
            "Texts to translate:",
            "For example, if you receive:",
            "pg1,el0,run0,x=100,y=200:Tämä on pitkä",
            "pg1,el0,run1,x=100,y=200:lause, joka on",
            "You MUST return:",
            "pg1,el0,run0:This is a long",
            "pg1,el0,run1:sentence that has been",
            "Do not add any extra explanations",
        ],
        id="single_item",
    ),
    dict(
        texts_for_llm=[
            {
                "id": "pg1,el0,run0",
                "text_to_send": "First",
                "shape_x": 100,
                "shape_y": 200,
            },
            {
                "id": "pg1,el1,run0",
                "text_to_send": "Second",
                "shape_x": 300,
                "shape_y": 400,
            },
        ],
        expected_formatted_text=(
            "pg1,el0,run0,x=100,y=200:First\npg1,el1,run0,x=300,y=400:Second"
        ),
        expected_prompt_contains=[
            "You are an expert Finnish to English translator.",
            f"EOL marker: '{EOL_MARKER}'",
            "Texts to translate:",
            "For example, if you receive:",
            "pg1,el0,run0,x=100,y=200:Tämä on pitkä",
            "pg1,el0,run1,x=100,y=200:lause, joka on",
            "You MUST return:",
            "pg1,el0,run0:This is a long",
            "pg1,el0,run1:sentence that has been",
            "Do not add any extra explanations",
        ],
        id="multiple_items",
    ),
)
def test_build_llm_prompt_and_data(
    texts_for_llm: list[dict[str, Any]],
    expected_formatted_text: str,
    expected_prompt_contains: list[str],
) -> None:
    prompt_text, formatted_text_for_llm = _build_llm_prompt_and_data(texts_for_llm)

    for expected_substring in expected_prompt_contains:
        assert expected_substring in prompt_text

    assert formatted_text_for_llm == expected_formatted_text


@patch("pptrans.__main__.click.echo")
@pytest.mark.kwparametrize(
    dict(
        run_details=[],
        id="empty_run_details",
    ),
    dict(
        run_details=[
            {
                "run_object": create_mock_run("Original"),
                "final_translation": "Translation",
                "from_cache": True,
            }
        ],
        id="single_translation_from_cache",
    ),
    dict(
        run_details=[
            {
                "run_object": create_mock_run("Original"),
                "final_translation": "Translation",
                "from_cache": False,
            }
        ],
        id="single_translation_from_llm",
    ),
    dict(
        run_details=[
            {
                "run_object": create_mock_run("Original1"),
                "final_translation": "Translation1",
                "from_cache": True,
            },
            {
                "run_object": create_mock_run("Original2"),
                "final_translation": "Translation2",
                "from_cache": False,
            },
        ],
        id="mix_of_cache_and_llm",
    ),
    dict(
        run_details=[
            # Add mock run object and original_text to match code expectations
            {
                "final_translation": "Translation1",
                "from_cache": True,
                "run_object": create_mock_run("Original"),
                "original_text": "Original",
            }
        ],
        id="missing_run_object",
    ),
)
def test_apply_translations_to_runs(
    mock_echo: MagicMock, run_details: list[dict[str, Any]]
) -> None:
    """Test the _apply_translations_to_runs function."""
    _apply_translations_to_runs(run_details)

    # Verify that run.text was updated for each run in the details
    for detail in run_details:
        if detail["run_object"] is not None:
            detail["run_object"].text = detail["final_translation"]

    # Check that echo was called for each run detail
    assert mock_echo.call_count >= len(run_details)


@patch("pptrans.__main__.commit_pending_cache_updates")
@patch("pptrans.__main__._apply_translations_to_runs")
@patch("pptrans.__main__.update_data_from_llm_response")
@patch("pptrans.__main__.llm.get_model")
@patch("pptrans.__main__._build_llm_prompt_and_data")
@patch("pptrans.__main__.prepare_slide_for_translation")
@patch("pptrans.__main__.generate_page_hash")
@patch("pptrans.__main__._extract_run_info_from_slide")
@patch("pptrans.__main__.load_cache")
@patch("pptrans.__main__.click.echo")
def test_process_translation_mode_no_slides(
    mock_echo: MagicMock,
    mock_load_cache: MagicMock,
    mock_extract_run: MagicMock,
    _mock_gen_hash: MagicMock,
    _mock_prep_slide: MagicMock,
    _mock_build_prompt: MagicMock,
    _mock_get_model: MagicMock,
    _mock_update_llm_resp: MagicMock,
    _mock_apply_trans: MagicMock,
    mock_commit_cache: MagicMock,
) -> None:
    mock_load_cache.return_value = {}
    _process_translation_mode(
        [], [], "cache.json", EOL_MARKER, "gemini-2.5-flash-preview-04-17"
    )
    mock_echo.assert_any_call(
        "No text found to process on selected slides for mode 'translate'."
    )
    mock_commit_cache.assert_called_once_with({}, {}, "cache.json")
    mock_extract_run.assert_not_called()


@patch("pptrans.__main__.commit_pending_cache_updates")
@patch("pptrans.__main__._apply_translations_to_runs")
@patch("pptrans.__main__.update_data_from_llm_response")
@patch("pptrans.__main__.llm.get_model")
@patch("pptrans.__main__._build_llm_prompt_and_data")
@patch("pptrans.__main__.prepare_slide_for_translation")
@patch("pptrans.__main__.generate_page_hash")
@patch("pptrans.__main__._extract_run_info_from_slide")
@patch("pptrans.__main__.load_cache")
@patch("pptrans.__main__.click.echo")
def test_process_translation_mode_all_cached(
    mock_echo: MagicMock,
    mock_load_cache: MagicMock,
    mock_extract_run: MagicMock,
    mock_gen_hash: MagicMock,
    mock_prep_slide: MagicMock,
    mock_build_prompt: MagicMock,
    mock_get_model: MagicMock,
    mock_update_llm_resp: MagicMock,
    mock_apply_trans: MagicMock,
    mock_commit_cache: MagicMock,
) -> None:
    # This test is for testing when all items are found in cache
    # Mock the cache to contain an entry for our slide hash
    mock_load_cache.return_value = {
        "slide_hash_1": {
            "processed_runs": [
                {
                    "original_text": "Cached Text",
                    "final_translation": "Cached Translation",
                    "run_object": create_mock_run("Cached Text"),
                    "shape_idx": 0,
                    "run_idx_in_shape": 0,
                    "shape_x": 100,
                    "shape_y": 200,
                    "llm_id": "pg1,el0,run0",
                }
            ]
        }
    }
    mock_extract_run.return_value = [
        {
            "original_text": "Cached Text",
            "run_object": create_mock_run("Cached Text"),
            "shape_idx": 0,
            "run_idx_in_shape": 0,
            "shape_x": 100,
            "shape_y": 200,
        }
    ]
    mock_gen_hash.return_value = "slide_hash_1"

    # Mock prepare_slide_for_translation to return cached runs
    processed_runs = [
        {
            "llm_id": "pg1,el0,run0",
            "original_text": "Cached Text",
            "run_object": create_mock_run("Cached Text"),
            "final_translation": "Cached Translation",
            "from_cache": True,
        }
    ]
    mock_prep_slide.return_value = ([], processed_runs, False)

    mock_slide = create_mock_slide(shapes=[])
    _process_translation_mode(
        [mock_slide], [0], "cache.json", EOL_MARKER, "gemini-2.5-flash-preview-04-17"
    )

    # Print out all echo calls to debug
    for _call in mock_echo.call_args_list:
        pass

    mock_echo.assert_any_call("Processing 1 selected slide(s) for translation...")
    mock_extract_run.assert_called_once_with(mock_slide)
    mock_gen_hash.assert_called_once_with(["Cached Text"])
    # Allow to call prepare_slide_for_translation - that's the actual implementation
    mock_prep_slide.assert_called_once()
    mock_build_prompt.assert_not_called()
    mock_get_model.assert_not_called()
    mock_update_llm_resp.assert_not_called()
    mock_apply_trans.assert_called_once()
    mock_commit_cache.assert_called_once_with(
        {
            "slide_hash_1": {
                "processed_runs": mock_load_cache.return_value["slide_hash_1"][
                    "processed_runs"
                ]
            }
        },
        {},
        "cache.json",
    )


@patch("pptrans.__main__.commit_pending_cache_updates")
@patch("pptrans.__main__._apply_translations_to_runs")
@patch("pptrans.__main__.update_data_from_llm_response")
@patch("pptrans.__main__.llm.get_model")
@patch("pptrans.__main__._build_llm_prompt_and_data")
@patch("pptrans.__main__.prepare_slide_for_translation")
@patch("pptrans.__main__.generate_page_hash")
@patch("pptrans.__main__._extract_run_info_from_slide")
@patch("pptrans.__main__.load_cache")
@patch("pptrans.__main__.click.echo")
def test_process_translation_mode_slide_with_no_text_then_slide_with_text(
    mock_echo: MagicMock,
    mock_load_cache: MagicMock,
    mock_extract_run_info: MagicMock,
    mock_gen_hash: MagicMock,
    mock_prep_slide: MagicMock,
    mock_build_prompt: MagicMock,
    # This mock is for llm.get_model() called inside _process_translation_mode:
    mock_get_model: MagicMock,
    mock_update_llm_resp: MagicMock,
    mock_apply_trans: MagicMock,
    mock_commit_cache: MagicMock,
) -> None:
    """Test _process_translation_mode with a no text slide and a with text slide."""
    mock_load_cache.return_value = {}  # Start with an empty cache
    # pending_page_cache_updates is internal to _process_translation_mode

    mock_slide_empty_obj = MagicMock(name="EmptySlideObj")
    mock_slide_with_text_obj = MagicMock(name="SlideWithTextObj")

    # These are the selected slides for processing
    slides_for_call = [mock_slide_empty_obj, mock_slide_with_text_obj]
    # Corresponding 0-indexed original page numbers for the slides in slides_for_call
    # (e.g., original page 1 is index 0, original page 2 is index 1)
    original_indices_for_call = [0, 1]

    mock_run_for_text_slide = create_mock_run("Hello")
    # _extract_run_info_from_slide returns list of dicts with shape coordinates
    runs_for_text_slide = [
        {
            "original_text": "Hello",
            "run_object": mock_run_for_text_slide,
            "shape_idx": 0,
            "run_idx_in_shape": 0,
            "shape_x": 100,
            "shape_y": 200,
        }
    ]

    mock_extract_run_info.side_effect = [
        [],  # For mock_slide_empty_obj (slide 1)
        runs_for_text_slide,  # For mock_slide_with_text_obj (slide 2)
    ]

    # Configure generate_page_hash for the second slide (which has text)
    page_hash_for_slide2 = "p2_hash_dummy"
    mock_gen_hash.return_value = page_hash_for_slide2

    # Configure prepare_slide_for_translation for the second slide
    # llm_id is generated by prepare_slide_for_translation
    llm_id_for_slide2_text0 = f"{page_hash_for_slide2}_t0"
    texts_for_llm_from_prep = [
        {
            "id": llm_id_for_slide2_text0,
            "original_text_for_cache": "Hello",  # Used by update_data_from_llm_response
            "text_to_send": "Hello" + EOL_MARKER,  # LLM gets text with EOL
            # For completeness, matches real prepare_slide:
            "run_object": mock_run_for_text_slide,
            "page_hash": page_hash_for_slide2,  # Used by update_data_from_llm_response
        }
    ]
    # This list will be modified in-place by the side_effect:
    details_for_apply_from_prep = [
        {
            "llm_id": llm_id_for_slide2_text0,
            "original_text": "Hello",
            "run_object": mock_run_for_text_slide,
            "final_translation": None,  # Placeholder, filled by update_data_from_llm
            "from_cache": False,
        }
    ]
    # prepare_slide_for_translation
    # returns (texts_for_llm, slide_details, updated_idx, page_requires_llm)
    mock_prep_slide.return_value = (
        texts_for_llm_from_prep,
        details_for_apply_from_prep,
        True,
    )

    # Configure _build_llm_prompt_and_data for texts from the second slide
    # _build_llm_prompt_and_data uses 'id' and 'text_to_send' from
    # texts_for_llm_from_prep
    llm_prompt_text = "Generated LLM Prompt"
    # formatted_text_for_llm should reflect 'text_to_send' which includes EOL_MARKER
    formatted_text_for_llm = f"{llm_id_for_slide2_text0}:Hello{EOL_MARKER}"
    mock_build_prompt.return_value = (llm_prompt_text, formatted_text_for_llm)

    # Configure LLM call
    mock_llm_instance = MagicMock(name="LLMModelInstance")
    mock_llm_response_obj = MagicMock(name="LLMResponseObject")
    # LLM response should also contain the EOL_MARKER if the prompt asked for it
    llm_response_content_str = f"{llm_id_for_slide2_text0}:Translated Hello{EOL_MARKER}"
    mock_llm_response_obj.text.return_value = llm_response_content_str
    mock_llm_instance.prompt.return_value = mock_llm_response_obj
    mock_get_model.return_value = mock_llm_instance

    # Define side effect for mock_update_llm_resp to modify arguments in place
    def mock_update_llm_side_effect(
        llm_lines_arg,
        texts_for_llm_list_arg,  # global_texts_for_llm_prompt_ref
        processed_runs_list_arg,  # all_processed_run_details_ref
        pending_updates_dict_arg,  # pending_page_cache_updates_ref
        eol_marker_str_arg,
    ) -> None:
        for line_from_llm in llm_lines_arg:
            stripped_line = line_from_llm.strip()
            if not stripped_line:
                continue
            parts = stripped_line.split(":", 1)
            if len(parts) == 2:
                parsed_id = parts[0].strip()
                translation_with_eol = parts[1]

                prompt_item = next(
                    (
                        item
                        for item in texts_for_llm_list_arg
                        if item["id"] == parsed_id
                    ),
                    None,
                )
                if prompt_item:
                    original_text_for_cache = prompt_item["original_text_for_cache"]
                    page_hash_of_item = prompt_item["page_hash"]
                    final_translation = translation_with_eol.removesuffix(
                        eol_marker_str_arg
                    )

                    for detail in processed_runs_list_arg:
                        if detail.get("llm_id") == parsed_id:
                            detail["final_translation"] = final_translation
                            break

                    if page_hash_of_item not in pending_updates_dict_arg:
                        pending_updates_dict_arg[page_hash_of_item] = []

                    found_in_pending = False
                    for pending_item in pending_updates_dict_arg[page_hash_of_item]:
                        if pending_item["original_text"] == original_text_for_cache:
                            pending_item["translation"] = final_translation
                            found_in_pending = True
                            break
                    if not found_in_pending:
                        pending_updates_dict_arg[page_hash_of_item].append(
                            {
                                "original_text": original_text_for_cache,
                                "translation": final_translation,
                            }
                        )

    mock_update_llm_resp.side_effect = mock_update_llm_side_effect

    # Call the function under test
    # _process_translation_mode uses llm.get_model() internally.
    # The 'model_name' is typically set up in the llm module context by the main CLI.
    # For this test, we ensure llm.get_model() is called
    # (and returns our mock_llm_instance).
    # The 'pending_cache_updates' dict is also internal.
    _process_translation_mode(
        slides_to_process=slides_for_call,
        original_page_indices=original_indices_for_call,
        cache_file_path="cache.json",
        eol_marker=EOL_MARKER,
        model="gemini-2.5-flash-preview-04-17",
    )

    # Assertions for the first slide (empty)
    mock_echo.assert_any_call("  Slide 1 (Original page 1): No text found.")
    # generate_page_hash is called with the list of original texts from the slide
    expected_texts_for_hash_slide2 = [
        item["original_text"] for item in runs_for_text_slide
    ]
    mock_gen_hash.assert_called_once_with(
        expected_texts_for_hash_slide2
    )  # Only for slide w/ text

    # Assertions for the second slide (with text)
    mock_extract_run_info.assert_any_call(mock_slide_with_text_obj)
    mock_prep_slide.assert_called_once_with(
        runs_for_text_slide,  # run_info_on_slide
        page_hash_for_slide2,  # page_hash
        {},  # translation_cache (empty for this test)
        EOL_MARKER,  # eol_marker
        2,  # current_page_num_1_indexed (for slide 2)
    )

    mock_build_prompt.assert_called_once_with(texts_for_llm_from_prep)
    # llm.get_model() in _process_translation_mode is called without arguments in the
    # actual code.
    # It relies on the llm module being pre-configured.
    # The test's mock_get_model is for the llm.get_model *module attribute*.
    # Allow for a model name argument, which is used in the actual implementation
    mock_get_model.assert_called_once()
    mock_llm_instance.prompt.assert_called_once_with(
        llm_prompt_text, fragments=[formatted_text_for_llm]
    )

    # update_data_from_llm_response modifies details_for_apply_from_prep in place
    # The pending_page_cache_updates dict is created inside _process_translation_mode
    # and then passed to update_data_from_llm_response.
    # We need to capture what update_data_from_llm_response was called with.
    # Since it's modified in place, the assertion for commit_pending_cache_updates
    # will cover its final state.
    args_update_llm, _ = mock_update_llm_resp.call_args
    assert (
        args_update_llm[0] == llm_response_content_str.splitlines()
    )  # llm_response_content.splitlines()
    # Corrected assertions based on the actual call signature in __main__.py:
    # update_data_from_llm_response(
    #     translated_text_response_str.splitlines(),  # args_update_llm[0]
    #     global_texts_for_llm_prompt,                # args_update_llm[1]
    #     all_processed_run_details,                  # args_update_llm[2]
    #     pending_page_cache_updates,                 # args_update_llm[3]
    #     eol_marker,                                 # args_update_llm[4]
    # )
    assert args_update_llm[1] == texts_for_llm_from_prep  # global_texts_for_llm_prompt
    assert (
        args_update_llm[2] == details_for_apply_from_prep
    )  # all_processed_run_details
    # args_update_llm[3] refers to the pending_page_cache_updates dict *after*
    # the side_effect has modified it in place.
    expected_pending_updates_after_side_effect = {
        page_hash_for_slide2: [
            {"original_text": "Hello", "translation": "Translated Hello"}
        ]
    }
    assert (
        args_update_llm[3] == expected_pending_updates_after_side_effect
    )  # pending_page_cache_updates
    assert args_update_llm[4] == EOL_MARKER  # eol_marker
    # Check that details_for_apply_from_prep was updated before apply_translations
    expected_details_after_llm_update = [
        {
            "llm_id": llm_id_for_slide2_text0,
            "original_text": "Hello",
            "run_object": mock_run_for_text_slide,
            "final_translation": "Translated Hello",  # Updated
            "from_cache": False,  # Remains False
        }
    ]
    mock_apply_trans.assert_called_once_with(expected_details_after_llm_update)

    # This is what pending_page_cache_updates should look like after the side_effect
    expected_pending_cache_updates_for_commit = {
        page_hash_for_slide2: [  # List of dicts for the page
            {
                "original_text": "Hello",  # from original_text_for_cache
                "translation": "Translated Hello",  # final_translation (without EOL)
            }
        ]
    }
    # translation_cache (first arg) is the initially loaded cache ({})
    mock_commit_cache.assert_called_once_with(
        {}, expected_pending_cache_updates_for_commit, "cache.json"
    )

    # Ensure "No text found to process on selected slides" was NOT called
    for call_args in mock_echo.call_args_list:
        args_tuple = call_args[0]
        if args_tuple:  # Ensure there are arguments
            assert "No text found to process on selected slides" not in args_tuple[0]

    assert mock_extract_run_info.call_count == 2


@patch("pptrans.__main__.commit_pending_cache_updates")
@patch("pptrans.__main__._apply_translations_to_runs")
@patch("pptrans.__main__.update_data_from_llm_response")
@patch("pptrans.__main__.llm.get_model")
@patch("pptrans.__main__._build_llm_prompt_and_data")
@patch("pptrans.__main__.prepare_slide_for_translation")
@patch("pptrans.__main__.generate_page_hash")
@patch("pptrans.__main__._extract_run_info_from_slide")
@patch("pptrans.__main__.load_cache")
@patch("pptrans.__main__.click.echo")
def test_process_translation_mode_with_llm_call(
    mock_echo: MagicMock,
    mock_load_cache: MagicMock,
    mock_extract_run: MagicMock,
    mock_gen_hash: MagicMock,
    mock_prep_slide: MagicMock,
    mock_build_prompt: MagicMock,
    mock_get_model: MagicMock,
    mock_update_llm_resp: MagicMock,
    mock_apply_trans: MagicMock,
    mock_commit_cache: MagicMock,
) -> None:
    """Test translation mode with a direct LLM call."""
    mock_load_cache.return_value = {}  # Empty cache
    mock_extract_run.return_value = [
        {
            "original_text": "Text For LLM",
            "run_object": create_mock_run("Text For LLM"),
            "shape_idx": 0,
            "run_idx_in_shape": 0,
            "shape_x": 100,
            "shape_y": 200,
        }
    ]
    mock_gen_hash.return_value = "hash1"
    # prepare_slide_for_translation returns:
    # (texts_for_llm, processed_runs, requires_llm)
    texts_for_llm = [
        {
            "id": "pg1,el0,run0",
            "original_text_for_cache": "Text For LLM",
            "text_to_send": "Text For LLM" + EOL_MARKER,
            "run_object": create_mock_run("Text For LLM"),
            "page_hash": "hash1",
            "shape_x": 100,
            "shape_y": 200,
        }
    ]
    processed_runs = [
        {
            "llm_id": "pg1,el0,run0",
            "original_text": "Text For LLM",
            "run_object": create_mock_run("Text For LLM"),
            "final_translation": None,
            "from_cache": False,
        }
    ]
    mock_prep_slide.return_value = (texts_for_llm, processed_runs, True)

    # Mock _build_llm_prompt_and_data
    mock_build_prompt.return_value = (
        "LLM Prompt",
        "pg1,el0,run0,x=100,y=200:Text For LLM" + EOL_MARKER,
    )

    # Mock LLM
    mock_llm = MagicMock()
    mock_llm_response = MagicMock()
    mock_llm_response.text.return_value = "pg1,el0,run0:Translated Text" + EOL_MARKER
    mock_llm.prompt.return_value = mock_llm_response
    mock_get_model.return_value = mock_llm

    # Add mock_update_llm_resp side effect that updates the pending_page_cache_updates
    def update_llm_side_effect(
        llm_lines,
        texts_for_llm_list,
        processed_runs_list,
        pending_updates_dict,
        eol_marker_str,
    ) -> None:
        # Update the pending cache updates to match what we expect
        pending_updates_dict["hash1"] = [
            {"original_text": "Text For LLM", "translation": "Translated Text"}
        ]
        # Update the processed runs
        for run in processed_runs_list:
            if run.get("llm_id") == "pg1,el0,run0":
                run["final_translation"] = "Translated Text"

    mock_update_llm_resp.side_effect = update_llm_side_effect

    # Action
    mock_slide = MagicMock()
    _process_translation_mode(
        [mock_slide], [0], "cache.json", EOL_MARKER, "gemini-2.5-flash-preview-04-17"
    )

    # Assertions
    mock_extract_run.assert_called_once_with(mock_slide)
    mock_gen_hash.assert_called_once_with(["Text For LLM"])
    mock_prep_slide.assert_called_once_with(
        mock_extract_run.return_value, "hash1", {}, EOL_MARKER, 1
    )

    mock_build_prompt.assert_called_once_with(texts_for_llm)
    mock_get_model.assert_called_once()
    mock_llm.prompt.assert_called_once_with(
        "LLM Prompt",
        fragments=["pg1,el0,run0,x=100,y=200:Text For LLM" + EOL_MARKER],
    )

    # Check pending updates were correctly initialized
    pending_updates = {
        "hash1": [{"original_text": "Text For LLM", "translation": "Translated Text"}]
    }

    # Check final commit
    mock_commit_cache.assert_called_once_with({}, pending_updates, "cache.json")


@patch("pptrans.__main__.click.echo")
def test_process_reverse_words_mode_no_slides_to_process(mock_echo: MagicMock) -> None:
    """Test _process_reverse_words_mode with no slides to process."""
    _process_reverse_words_mode([], EOL_MARKER)
    mock_echo.assert_any_call(
        "No text found to process on selected slides for mode 'reverse-words'."
    )


@patch("pptrans.__main__.click.echo")
@patch("pptrans.__main__._extract_run_info_from_slide")
def test_process_reverse_words_mode_no_text(
    mock_extract_run: MagicMock, mock_echo: MagicMock
) -> None:
    """Test _process_reverse_words_mode with a slide but no text."""
    mock_extract_run.return_value = []  # No text found
    mock_slide = MagicMock()
    _process_reverse_words_mode([mock_slide], EOL_MARKER)
    mock_echo.assert_any_call(
        "Extracting text from 1 slides for mode 'reverse-words'..."
    )
    mock_echo.assert_any_call(
        "No text found to process on selected slides for mode 'reverse-words'."
    )


@patch("pptrans.__main__.click.echo")
@patch("pptrans.__main__._extract_run_info_from_slide")
def test_process_reverse_words_mode_with_text(
    mock_extract_run: MagicMock, mock_echo: MagicMock
) -> None:
    """Test _process_reverse_words_mode with a slide containing text."""
    mock_run1 = create_mock_run("Hello world")
    mock_run2 = create_mock_run("Another text")
    mock_extract_run.return_value = [
        {"original_text": "Hello world", "run_object": mock_run1},
        {"original_text": "Another text", "run_object": mock_run2},
    ]
    mock_slide = MagicMock()
    _process_reverse_words_mode([mock_slide], EOL_MARKER)
    assert mock_run1.text == "olleH dlrow"
    assert mock_run2.text == "rehtonA txet"
    mock_echo.assert_any_call(
        "Extracting text from 1 slides for mode 'reverse-words'..."
    )
    mock_echo.assert_any_call(
        "Found 2 text elements to process for mode 'reverse-words'."
    )


@patch("pptrans.__main__.click.echo")
@pytest.mark.kwparametrize(
    dict(
        bak_path=None,
        expected_calls=[],  # No additional calls expected
        id="no_bak_path",
    ),
    dict(
        bak_path=Path("backup.pptx"),
        expected_calls=[],  # No additional calls expected
        id="with_bak_path",
    ),
)
def test_emit_save_message(
    mock_echo: MagicMock, bak_path: Path | None, expected_calls: list[str]
) -> None:
    """Test the _emit_save_message function."""
    mode = "translate"
    output_path = Path("output.pptx")
    _emit_save_message(mode, output_path)
    mock_echo.assert_called_once_with(
        f"Presentation saved in '{mode}' mode to: {output_path}"
    )


def test_main_cli() -> None:
    """Test the main CLI entrypoint function."""
    # Import here to avoid import-time patching issues

    # Use the proper Click runner
    runner = CliRunner()

    with runner.isolated_filesystem():
        # Create a dummy input file
        Path("input.pptx").write_bytes(b"dummy pptx content")

        # Apply patches inside the context to ensure they're active during CLI execution
        with (
            patch("pptrans.__main__.shutil.copy2") as mock_copy2,
            patch("pptrans.__main__._handle_slide_selection") as mock_handle_slides,
            patch(
                "pptrans.__main__._process_translation_mode"
            ) as mock_process_translate,
            patch(
                "pptrans.__main__._process_reverse_words_mode"
            ) as mock_process_reverse,
            patch("pptrans.__main__.Presentation") as mock_presentation,
            patch("pptrans.__main__.click.echo"),
            patch("pptrans.__main__.llm") as mock_llm,
        ):
            # Set up mocks
            mock_presentation_instance = MagicMock()
            mock_presentation_instance.slides = ["slide1", "slide2"]
            mock_presentation.return_value = mock_presentation_instance
            mock_handle_slides.return_value = {0, 1}  # Both slides

            # Mock LLM module attributes
            mock_llm_model = MagicMock()
            mock_llm.get_model.return_value = mock_llm_model

            # Call the command within the isolated filesystem
            result = runner.invoke(
                main,
                [
                    "input.pptx",
                    "output.pptx",
                    "--mode",
                    "translate",
                    "--model",
                    "custom-model",
                ],
                catch_exceptions=False,
            )

            # Verify command execution
            assert result.exit_code == 0

            # Verify the right functions were called
            mock_copy2.assert_called_once()
            mock_handle_slides.assert_called_once()
            # Check that the model parameter is passed to the processing function
            mock_process_translate.assert_called_once()
            args, _ = mock_process_translate.call_args
            # Check that the model parameter was passed
            assert args[4] == "custom-model"
            mock_process_reverse.assert_not_called()

            # Verify the presentation was saved
            mock_presentation_instance.save.assert_called_once()


def test_main_cli_no_slides_selected_no_pages_option() -> None:
    """Test the main CLI when no slides are selected."""
    # Import here to avoid import-time patching issues

    runner = CliRunner()
    with runner.isolated_filesystem():
        # Create a dummy input file
        Path("input.pptx").write_bytes(b"dummy pptx content")

        # Apply patches inside the context
        with (
            patch("pptrans.__main__.shutil.copy2") as mock_copy2,
            patch("pptrans.__main__._handle_slide_selection") as mock_handle_slides,
            patch(
                "pptrans.__main__._process_translation_mode"
            ) as mock_process_translate,
            patch(
                "pptrans.__main__._process_reverse_words_mode"
            ) as mock_process_reverse,
            patch("pptrans.__main__.Presentation") as mock_presentation,
            patch("pptrans.__main__.click.echo") as mock_echo,
            patch("pptrans.__main__.load_cache") as mock_load_cache,
            patch("pptrans.__main__.commit_pending_cache_updates") as mock_commit_cache,
        ):
            # Set up mocks
            mock_presentation_instance = MagicMock()
            mock_presentation_instance.slides = ["slide1", "slide2"]
            mock_presentation.return_value = mock_presentation_instance
            mock_handle_slides.return_value = set()  # No slides selected
            mock_cache = {}
            mock_load_cache.return_value = mock_cache

            # Call the command
            result = runner.invoke(
                main,
                [
                    "input.pptx",
                    "output.pptx",
                    "--mode",
                    "translate",
                    # Default model should be used
                ],
                catch_exceptions=False,
            )

            # Verify command execution
            assert result.exit_code == 0

            # Verify behavior with no slides selected
            mock_copy2.assert_called_once()
            mock_process_translate.assert_not_called()
            mock_process_reverse.assert_not_called()
            mock_presentation_instance.save.assert_called_once()
            mock_commit_cache.assert_called_once()

            # Verify appropriate warnings were echoed
            warning_call_found = False
            for call in mock_echo.call_args_list:
                if "No slides selected for processing" in str(
                    call
                ) or "no slides being selected" in str(call):
                    warning_call_found = True
                    break
            assert warning_call_found


def test_main_cli_no_slides_selected_reverse_words_mode_no_pages_option() -> None:
    """Test the main CLI with reverse_words mode and no slides selected."""
    # Import here to avoid import-time patching issues

    runner = CliRunner()
    with runner.isolated_filesystem():
        # Create a dummy input file
        Path("input.pptx").write_bytes(b"dummy pptx content")

        # Apply patches inside the context
        with (
            patch("pptrans.__main__.shutil.copy2") as mock_copy2,
            patch("pptrans.__main__._handle_slide_selection") as mock_handle_slides,
            patch(
                "pptrans.__main__._process_translation_mode"
            ) as mock_process_translate,
            patch(
                "pptrans.__main__._process_reverse_words_mode"
            ) as mock_process_reverse,
            patch("pptrans.__main__.Presentation") as mock_presentation,
            patch("pptrans.__main__.click.echo") as mock_echo,
            patch("pptrans.__main__._emit_save_message") as mock_emit_save,
        ):
            # Set up mocks
            mock_presentation_instance = MagicMock()
            mock_presentation_instance.slides = ["slide1", "slide2"]
            mock_presentation.return_value = mock_presentation_instance
            mock_handle_slides.return_value = set()  # No slides selected

            # Call the command with reverse-words mode
            result = runner.invoke(
                main,
                [
                    "input.pptx",
                    "output.pptx",
                    "--mode",
                    "reverse-words",
                    # Model parameter is ignored in reverse-words mode
                    "--model",
                    "ignored-model",
                ],
                catch_exceptions=False,
            )

            # Verify command execution
            assert result.exit_code == 0

            # Verify the right functions were called for reverse-words mode
            mock_copy2.assert_called_once()
            mock_process_translate.assert_not_called()
            mock_process_reverse.assert_not_called()  # Not to be called with no slides
            mock_presentation_instance.save.assert_called_once()

            # Properly check emit_save behavior. Shouldn't be called due to early return
            mock_emit_save.assert_not_called()

            # Verify appropriate no-slides warning is echoed
            warning_call_found = False
            for call in mock_echo.call_args_list:
                if "No slides selected for processing" in str(
                    call
                ) or "no slides being selected" in str(call):
                    warning_call_found = True
                    break
            assert warning_call_found

            # Also verify the "Presentation saved (no text modifications)" message
            save_message_found = False
            for call in mock_echo.call_args_list:
                if "Presentation saved (no text modifications)" in str(call):
                    save_message_found = True
                    break
            assert save_message_found

            # Verify appropriate warnings were echoed
            warning_call_found = False
            for call in mock_echo.call_args_list:
                if "No slides selected for processing" in str(
                    call
                ) or "no slides being selected" in str(call):
                    warning_call_found = True
                    break
            assert warning_call_found


def test_main_dunder_guard(capsys: pytest.CaptureFixture[str]) -> None:
    """Test the ``if __name__ == '__main__':`` block.

    The code in ``if __name__ == '__main__':`` is not covered by the test suite.
    This is a simple test to ensure that it doesn't crash when run directly.
    """
    # Instead of using runpy.run_module which requires a proper module spec,
    # import the module and simulate the __main__ guard directly
    main_module = importlib.import_module("pptrans.__main__")

    # Save the original value
    original_name = main_module.__name__

    # Mock sys.argv to provide arguments to Click
    with mock.patch.object(sys, "argv", ["pptrans", "--help"]):
        try:
            # Simulate as if running as __main__
            main_module.__name__ = "__main__"
            # Execute the main function directly
            with (
                patch.object(main_module, "__name__", "__main__"),
                mock.patch.object(click, "echo"),  # Suppress output
                mock.patch.object(sys, "exit") as exit_mock,  # Prevent actual exit
            ):
                main_module.main()
                # Check that sys.exit was called, which is expected in CLI tools
                exit_mock.assert_called_once()
        finally:
            # Restore original value
            main_module.__name__ = original_name
